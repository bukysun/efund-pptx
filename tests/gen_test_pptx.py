"""
生成用于 QA 测试的示例 PPTX。
覆盖封面、目录、正文（两节内容）、封底四种布局。
刻意在第三张幻灯片正文中留一个 XXXX 占位符，用于验证内容检查能否检出。

用法:
    python3.12 tests/gen_test_pptx.py [output_path]
    默认输出到 tests/output_test.pptx
"""

import copy
import math
import os
import sys

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE   = os.path.join(SKILL_DIR, "assets", "template.pptx")

# ── 颜色 ──────────────────────────────────────────────────────
DEEP_BLUE   = RGBColor(0,   80,  150)
BRIGHT_BLUE = RGBColor(30,  185, 225)
DARK_GRAY   = RGBColor(60,  60,  60)
WHITE       = RGBColor(255, 255, 255)
LIGHT_GRAY  = RGBColor(204, 204, 204)
PALE_GRAY   = RGBColor(242, 242, 242)

CN_FONT = "华文黑体_易方达"
EN_FONT = "Arial"

# ── 内容区常量 ────────────────────────────────────────────────
_PH10_LEFT  = Inches(0.40)
_PH10_TOP   = Inches(1.42)
_PH10_WIDTH = {2: Inches(8.782), 3: Inches(4.884)}
CONTENT_HEIGHT_PT = 2.71 * 72


# ── 字体工具 ──────────────────────────────────────────────────

def apply_font(run, cn_font=CN_FONT, en_font=EN_FONT,
               size=None, bold=None, color=None, italic=None, underline=None):
    rPr = run._r.get_or_add_rPr()
    if cn_font and not en_font:
        for tag, face in [('a:latin', cn_font), ('a:cs', cn_font)]:
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set('typeface', face)
            el.set('charset', '-122')
    elif en_font and not cn_font:
        latin = rPr.find(qn('a:latin'))
        if latin is None:
            latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', en_font)
        latin.attrib.pop('charset', None)
    elif cn_font and en_font:
        for tag, face in [('a:latin', en_font), ('a:ea', cn_font)]:
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set('typeface', face)
    if size      is not None: run.font.size      = Pt(size)
    if bold      is not None: run.font.bold      = bold
    if italic    is not None: run.font.italic    = italic
    if underline is not None: run.font.underline = underline
    if color     is not None: run.font.color.rgb = color


def add_text(tf, text, *, first=False, align=PP_ALIGN.LEFT,
             cn_font=CN_FONT, en_font=EN_FONT,
             size=None, bold=None, color=DEEP_BLUE,
             italic=None, underline=None):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    if first:
        para.clear()
    run = para.add_run()
    run.text = text
    apply_font(run, cn_font=cn_font, en_font=en_font,
               size=size, bold=bold, color=color, italic=italic, underline=underline)
    para.alignment = align
    return para


# ── 标题注入 ──────────────────────────────────────────────────

_TITLE_SP_XML = '''<p:sp
  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="9999" name="TitleInjected"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="title"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="323850" y="211138"/><a:ext cx="8382375" cy="268287"/></a:xfrm>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/><a:lstStyle/>
    <a:p><a:r><a:t></a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''


def set_slide_title(slide, text, size=23, bold=True, color=WHITE):
    title_shape = slide.shapes.title
    if title_shape is None:
        slide.shapes._spTree.insert(2, etree.fromstring(_TITLE_SP_XML))
        title_shape = slide.shapes.title
    if title_shape:
        if title_shape.width < Inches(8):
            _l = title_shape.left; _t = title_shape.top; _h = title_shape.height
            title_shape.left = _l; title_shape.top = _t
            title_shape.height = _h; title_shape.width = Inches(9.17)
        tf = title_shape.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        add_text(tf, text, first=True,
                 cn_font=CN_FONT, en_font=None,
                 size=size, bold=bold, color=color)


# ── 正文工具 ──────────────────────────────────────────────────

BODY_STYLES = {
    0: dict(size=15, bold=True,  color=DEEP_BLUE, cn_font=CN_FONT, en_font=EN_FONT),
    1: dict(size=12, bold=False, color=DARK_GRAY,  cn_font=CN_FONT, en_font=EN_FONT),
    2: dict(size=10, bold=False, color=DARK_GRAY,  cn_font=CN_FONT, en_font=EN_FONT),
}

_LINE_HT_PT = {0: 15*1.40, 1: 12*1.55, 2: 10*1.55}
_BULLET_INDENT_EMU = 171450


def _set_para_spacing(para, spc_before_pt=None, line_spc_pct=None):
    pPr = para._p.get_or_add_pPr()
    if line_spc_pct is not None:
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        lnSpc.clear()
        etree.SubElement(lnSpc, qn('a:spcPct')).set('val', str(int(line_spc_pct * 1000)))
    if spc_before_pt is not None:
        spcBef = pPr.find(qn('a:spcBef'))
        if spcBef is None:
            spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcBef.clear()
        etree.SubElement(spcBef, qn('a:spcPts')).set('val', str(int(spc_before_pt * 100)))


def _set_para_bullet(para, enabled=True, level=1):
    pPr = para._p.get_or_add_pPr()
    for tag in [qn('a:buNone'), qn('a:buChar'), qn('a:buAutoNum'),
                qn('a:buFont'), qn('a:buSzPct')]:
        el = pPr.find(tag)
        if el is not None:
            pPr.remove(el)
    if not enabled:
        etree.SubElement(pPr, qn('a:buNone'))
        pPr.attrib.pop('marL', None)
        pPr.attrib.pop('indent', None)
        return
    pPr.set('marL', str(_BULLET_INDENT_EMU * level))
    pPr.set('indent', str(-_BULLET_INDENT_EMU))
    buFont = etree.SubElement(pPr, qn('a:buFont'))
    buFont.set('typeface', 'Arial')
    buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
    buSzPct.set('val', '100000')
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', '•')


def add_body_content(tf, items, available_pt=CONTENT_HEIGHT_PT, available_width_pt=None):
    _MIN_L0_SPC = 6.0
    _TARGET_LNS = {0: 140, 1: 160, 2: 160}
    _FONT_SZ    = {0: 15,  1: 12,  2: 10}

    n_body_items = sum(1 for _, lv in items if lv != 0)
    use_bullet   = (n_body_items > 1)

    def _est_lines(text, level):
        if available_width_pt is None:
            return 1
        font_sz   = _FONT_SZ.get(level, 12)
        indent_pt = (_BULLET_INDENT_EMU / 12700) if (use_bullet and level > 0) else 0
        usable    = max(font_sz, available_width_pt - indent_pt - 18)
        chars_line = max(1, usable / font_sz)
        return max(1, math.ceil(len(text) / chars_line))

    n_l0_gaps  = sum(1 for i, (_, lv) in enumerate(items) if lv == 0 and i > 0)
    avail_text = available_pt - n_l0_gaps * _MIN_L0_SPC
    base_target = sum(_FONT_SZ.get(lv,12) * _TARGET_LNS.get(lv,160)/100
                      * _est_lines(t, lv) for t, lv in items)
    base_min    = sum(_FONT_SZ.get(lv,12) * _est_lines(t, lv) for t, lv in items)

    if base_target <= avail_text:
        scale     = 1.0
        remaining = avail_text - base_target
    elif base_min <= avail_text:
        scale     = avail_text / base_target
        remaining = 0.0
    else:
        scale     = 100 / max(_TARGET_LNS.values())
        remaining = 0.0

    actual_lns = {lv: max(100, int(_TARGET_LNS[lv] * scale)) for lv in [0, 1, 2]}

    _MAX_L0_EXTRA = 28.0
    _MAX_OTHER    = 12.0
    n_other_gaps  = sum(1 for idx, (_, lv) in enumerate(items) if idx > 0 and lv != 0)

    if n_l0_gaps > 0 and remaining > 0:
        total_w  = 4 * n_l0_gaps + n_other_gaps
        unit_raw = remaining / total_w
        l0_extra = min(_MAX_L0_EXTRA, 4 * unit_raw)
    else:
        l0_extra = 0.0

    other_budget = remaining - n_l0_gaps * l0_extra
    unit_other   = (other_budget / n_other_gaps) if n_other_gaps > 0 else 0.0
    unit_other   = min(_MAX_OTHER, unit_other)

    for i, (text, level) in enumerate(items):
        para = add_text(tf, text, first=(i == 0), **BODY_STYLES.get(level, BODY_STYLES[1]))
        if level == 0 or not use_bullet:
            _set_para_bullet(para, enabled=False)
        else:
            _set_para_bullet(para, enabled=True, level=level)
        if i == 0:
            spc_before = 0.0
        elif level == 0:
            spc_before = _MIN_L0_SPC + l0_extra
        else:
            spc_before = unit_other
        _set_para_spacing(para, spc_before_pt=spc_before, line_spc_pct=actual_lns[level])


# ── 目录工具 ──────────────────────────────────────────────────

def _set_toc_row(tr, number, title, color_hex):
    cells = tr.findall(qn('a:tc'))

    def _apply_color(rEl, hex_val):
        rPr = rEl.find(qn('a:rPr'))
        if rPr is None:
            rPr = etree.SubElement(rEl, qn('a:rPr'))
        old = rPr.find(qn('a:solidFill'))
        if old is not None:
            rPr.remove(old)
        fill = etree.Element(qn('a:solidFill'))
        etree.SubElement(fill, qn('a:srgbClr')).set('val', hex_val)
        rPr.insert(0, fill)

    num_color = color_hex if color_hex == 'CCCCCC' else '1EB9E1'

    for rEl in cells[0].findall('.//' + qn('a:r')):
        t = rEl.find(qn('a:t'))
        if t is not None:
            t.text = number
        _apply_color(rEl, num_color)

    for rEl in cells[1].findall('.//' + qn('a:r')):
        t = rEl.find(qn('a:t'))
        if t is not None:
            t.text = title
        _apply_color(rEl, color_hex)


def fill_toc_table(tbl, chapters, active_idx):
    tbl_xml = tbl._tbl
    tr_list = tbl_xml.findall(qn('a:tr'))
    n_have  = len(tr_list)
    n_need  = len(chapters)
    if n_need < n_have:
        for tr in tr_list[n_need:]:
            tbl_xml.remove(tr)
    elif n_need > n_have:
        style_row = tr_list[-1]
        for _ in range(n_need - n_have):
            tbl_xml.append(copy.deepcopy(style_row))
    tr_list = tbl_xml.findall(qn('a:tr'))
    for ri, tr in enumerate(tr_list):
        if active_idx is None:
            color_hex = '005096'
        else:
            color_hex = '005096' if ri == active_idx else 'CCCCCC'
        _set_toc_row(tr, f'{ri + 1:02d}.', chapters[ri], color_hex)


def add_toc_slide(prs, toc_table_xml, chapters, active_idx=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    cloned = copy.deepcopy(toc_table_xml)
    slide.shapes._spTree.append(cloned)
    for shape in slide.shapes:
        if shape.shape_type == 19:
            fill_toc_table(shape.table, chapters, active_idx)
            break
    return slide


# ── 封底工具 ──────────────────────────────────────────────────

def add_cn_closing_slide(prs, thanks='谢谢', contacts=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        tf  = ph.text_frame
        tf.word_wrap = True
        if idx == 10:
            para = add_text(tf, thanks, first=True,
                            cn_font=CN_FONT, en_font=None,
                            size=45, bold=True, color=DEEP_BLUE)
            _set_para_spacing(para, line_spc_pct=100)
        elif idx == 11 and contacts:
            for i, line in enumerate(contacts):
                para = add_text(tf, line, first=(i == 0),
                                cn_font=CN_FONT, en_font=None,
                                size=12, bold=False, color=DEEP_BLUE)
                _set_para_spacing(para, line_spc_pct=170)
    return slide


# ── 免责声明 ──────────────────────────────────────────────────

def load_disclaimer(lang="cn"):
    path = os.path.join(SKILL_DIR, "disclaimers", f"disclaimer_{lang}.txt")
    with open(path, encoding="utf-8") as f:
        return f.read()


# ═══════════════════════════════════════════════════════════════
# 主程序
# ═══════════════════════════════════════════════════════════════

def build(output_path):
    prs = Presentation(TEMPLATE)

    # 提取目录表格 XML（必须在删除 slides 前）
    toc_table_xml = None
    for slide in prs.slides:
        if slide.slide_layout.name == '中文目录页，仅供目录页使用':
            for shape in slide.shapes:
                if shape.shape_type == 19:
                    toc_table_xml = copy.deepcopy(shape._element)
                    break
            break

    # 删除所有模板 slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    chapters = ["公司简介", "投资策略", "产品线概览"]
    disclaimer = load_disclaimer("cn")

    # ── 1. 中文封面 ────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in slide.shapes:
        if not hasattr(shape, 'placeholder_format'):
            continue
        idx = shape.placeholder_format.idx
        tf  = shape.text_frame
        if idx == 1:
            add_text(tf, "易方达基金管理有限公司", first=True, size=28, bold=False)
            add_text(tf, "2024年度报告", size=22, bold=True)
        elif idx == 10:
            add_text(tf, "汇报人：张三", first=True, size=14, bold=False)
            add_text(tf, "2024年3月", size=14, bold=False)

    # ── 2. 总览目录页 ──────────────────────────────────────────
    add_toc_slide(prs, toc_table_xml, chapters)

    # ── 3. 正文页（含刻意留下的 XXXX 占位符，用于测试检测） ──
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_slide_title(slide, '公司简介')
    ph10 = None
    for shape in slide.shapes:
        try:
            if shape.placeholder_format.idx == 10:
                ph10 = shape
                break
        except Exception:
            pass
    if ph10:
        add_body_content(ph10.text_frame, [
            ('规模概况', 0),
            ('截至2024年底，管理规模突破XXXX亿元', 1),   # ← 刻意留下的占位符
            ('公募基金规模行业前三', 1),
            ('主要产品线', 0),
            ('权益类基金：占比约40%', 1),
            ('固收类基金：占比约45%', 1),
        ], available_pt=CONTENT_HEIGHT_PT,
           available_width_pt=_PH10_WIDTH[2] / 12700)

    # ── 4. 正文页（章节跳转页 + 正常内容） ────────────────────
    add_toc_slide(prs, toc_table_xml, chapters, active_idx=1)

    slide2 = prs.slides.add_slide(prs.slide_layouts[2])
    set_slide_title(slide2, '投资策略')
    ph10b = None
    for shape in slide2.shapes:
        try:
            if shape.placeholder_format.idx == 10:
                ph10b = shape
                break
        except Exception:
            pass
    if ph10b:
        add_body_content(ph10b.text_frame, [
            ('核心理念', 0),
            ('以基本面研究驱动投资决策，注重长期价值创造', 1),
            ('风险收益比优先，坚持价值导向投资原则', 1),
            ('风险管理', 0),
            ('全流程风险管理，覆盖投前、投中、投后各环节', 1),
            ('独立风控团队，与投资团队形成制衡', 1),
        ], available_pt=CONTENT_HEIGHT_PT,
           available_width_pt=_PH10_WIDTH[2] / 12700)

    # ── 5. 封底 ────────────────────────────────────────────────
    add_cn_closing_slide(prs,
        thanks='谢谢',
        contacts=['联系我们：investor@efunds.com.cn']
    )

    prs.save(output_path)
    print(f"生成完成: {output_path}")


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "output_test.pptx"
    )
    build(out)
