"""End-to-end test: 完整 PPT 从封面到封底。"""
import sys, os, copy
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from lxml import etree

SKILL_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE  = os.path.join(SKILL_DIR, "assets", "template.pptx")

CN_FONT     = "华文黑体_易方达"
EN_FONT     = "Arial"
DEEP_BLUE   = RGBColor(0,   80,  150)
BRIGHT_BLUE = RGBColor(30,  185, 225)
DARK_GRAY   = RGBColor(60,  60,   60)
MID_GRAY    = RGBColor(150, 150, 150)
LIGHT_GRAY  = RGBColor(204, 204, 204)
WHITE       = RGBColor(255, 255, 255)

# ── 基础工具 ─────────────────────────────────────────────────

def _qn(tag):
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    prefix, local = tag.split(':')
    return f'{{{nsmap[prefix]}}}{local}'


def apply_font(run, cn_font=CN_FONT, en_font=None, size=None, bold=None, color=None):
    rPr = run._r.get_or_add_rPr()
    if cn_font and not en_font:
        latin = rPr.find(_qn('a:latin'))
        if latin is None: latin = etree.SubElement(rPr, _qn('a:latin'))
        latin.set('typeface', cn_font); latin.set('charset', '-122')
        cs = rPr.find(_qn('a:cs'))
        if cs is None: cs = etree.SubElement(rPr, _qn('a:cs'))
        cs.set('typeface', cn_font); cs.set('charset', '-122')
    elif en_font and not cn_font:
        latin = rPr.find(_qn('a:latin'))
        if latin is None: latin = etree.SubElement(rPr, _qn('a:latin'))
        latin.set('typeface', en_font); latin.attrib.pop('charset', None)
    elif cn_font and en_font:
        for tag, face in [('a:latin', en_font), ('a:ea', cn_font)]:
            el = rPr.find(_qn(tag))
            if el is None: el = etree.SubElement(rPr, _qn(tag))
            el.set('typeface', face)
    if size  is not None: run.font.size      = Pt(size)
    if bold  is not None: run.font.bold      = bold
    if color is not None: run.font.color.rgb = color


def add_text(tf, text, first=False, cn_font=CN_FONT, en_font=None,
             size=None, bold=None, color=DEEP_BLUE):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    if first: para.clear()
    run = para.add_run(); run.text = text
    apply_font(run, cn_font=cn_font, en_font=en_font, size=size, bold=bold, color=color)
    return para


def _set_para_spacing(para, spc_before_pt=None, line_spc_pct=None):
    pPr = para._p.get_or_add_pPr()
    if line_spc_pct is not None:
        lnSpc = pPr.find(_qn('a:lnSpc'))
        if lnSpc is None: lnSpc = etree.SubElement(pPr, _qn('a:lnSpc'))
        lnSpc.clear()
        etree.SubElement(lnSpc, _qn('a:spcPct')).set('val', str(int(line_spc_pct * 1000)))
    if spc_before_pt is not None:
        spcBef = pPr.find(_qn('a:spcBef'))
        if spcBef is None: spcBef = etree.SubElement(pPr, _qn('a:spcBef'))
        spcBef.clear()
        etree.SubElement(spcBef, _qn('a:spcPts')).set('val', str(int(spc_before_pt * 100)))

# ── 标题注入 ─────────────────────────────────────────────────

_TITLE_SP_XML = '''<p:sp
  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="9999" name="TitleInjected"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="title"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="323850" y="211138"/><a:ext cx="8382375" cy="268287"/></a:xfrm>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/><a:lstStyle/>
    <a:p><a:r><a:t></a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''


def set_slide_title(slide, text, size=23, bold=True, color=WHITE):
    title_shape = slide.shapes.title
    if title_shape is None:
        slide.shapes._spTree.insert(2, etree.fromstring(_TITLE_SP_XML))
        title_shape = slide.shapes.title
    if title_shape:
        add_text(title_shape.text_frame, text, first=True,
                 cn_font=CN_FONT, en_font=None, size=size, bold=bold, color=color)

# ── 正文 ─────────────────────────────────────────────────────

BODY_STYLES = {
    0: dict(size=15, bold=True,  color=DEEP_BLUE, cn_font=CN_FONT, en_font=None),
    1: dict(size=12, bold=False, color=DARK_GRAY,  cn_font=CN_FONT, en_font=None),
    2: dict(size=10, bold=False, color=DARK_GRAY,  cn_font=CN_FONT, en_font=None),
}
CONTENT_HEIGHT_PT = 2.71 * 72


def add_body_content(tf, items, available_pt=CONTENT_HEIGHT_PT):
    _MIN_L0_SPC = 6.0
    _TARGET_LNS = {0: 140, 1: 160, 2: 160}
    _FONT_SZ    = {0: 15,  1: 12,  2: 10}
    n_l0_gaps  = sum(1 for i, (_, lv) in enumerate(items) if lv == 0 and i > 0)
    avail_text = available_pt - n_l0_gaps * _MIN_L0_SPC
    base_target = sum(_FONT_SZ.get(lv,12) * _TARGET_LNS.get(lv,160)/100 for _,lv in items)
    base_min    = sum(_FONT_SZ.get(lv,12) for _,lv in items)
    if base_target <= avail_text:
        scale = 1.0; remaining = avail_text - base_target
    elif base_min <= avail_text:
        scale = avail_text / base_target; remaining = 0.0
    else:
        scale = 100 / max(_TARGET_LNS.values()); remaining = 0.0
    actual_lns = {lv: max(100, int(_TARGET_LNS[lv] * scale)) for lv in [0,1,2]}
    weights = [0 if i == 0 else (3 if lv == 0 else 1) for i, (_, lv) in enumerate(items)]
    total_w = sum(weights); unit = remaining / total_w if total_w > 0 else 0
    for i, (text, level) in enumerate(items):
        para = add_text(tf, text, first=(i == 0), **BODY_STYLES.get(level, BODY_STYLES[1]))
        if i == 0: spc_before = 0.0
        elif level == 0: spc_before = _MIN_L0_SPC + min(20.0 - _MIN_L0_SPC, 3 * unit)
        else: spc_before = min(20.0, unit)
        _set_para_spacing(para, spc_before_pt=spc_before, line_spc_pct=actual_lns[level])

# ── 目录 ─────────────────────────────────────────────────────

def _set_toc_row(tr, number, title, color_hex):
    cells = tr.findall(qn('a:tc'))
    for rEl in cells[0].findall('.//' + qn('a:r')):
        t = rEl.find(qn('a:t'))
        if t is not None: t.text = number
    for rEl in cells[1].findall('.//' + qn('a:r')):
        t = rEl.find(qn('a:t'))
        if t is not None: t.text = title
        rPr = rEl.find(qn('a:rPr'))
        if rPr is None: rPr = etree.SubElement(rEl, qn('a:rPr'))
        old = rPr.find(qn('a:solidFill'))
        if old is not None: rPr.remove(old)
        fill = etree.Element(qn('a:solidFill'))
        srgb = etree.SubElement(fill, qn('a:srgbClr'))
        srgb.set('val', color_hex)
        rPr.insert(0, fill)


def add_toc_slide(prs, toc_table_xml, chapters, active_idx):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    cloned = copy.deepcopy(toc_table_xml)
    slide.shapes._spTree.append(cloned)
    for shape in slide.shapes:
        if shape.shape_type == 19:
            tbl_xml = shape.table._tbl
            tr_list = tbl_xml.findall(qn('a:tr'))
            n_have, n_need = len(tr_list), len(chapters)
            if n_need < n_have:
                for tr in tr_list[n_need:]: tbl_xml.remove(tr)
            elif n_need > n_have:
                style_row = tr_list[-1]
                for _ in range(n_need - n_have): tbl_xml.append(copy.deepcopy(style_row))
            tr_list = tbl_xml.findall(qn('a:tr'))
            for ri, tr in enumerate(tr_list):
                color_hex = '005096' if ri == active_idx else 'CCCCCC'
                _set_toc_row(tr, f'{ri+1:02d}.', chapters[ri], color_hex)
            break
    return slide

# ── Layout 4 ─────────────────────────────────────────────────

_L4_LEFT_CHART  = (Inches(0.40), Inches(1.49), Inches(4.32), Inches(2.54))
_L4_RIGHT_CHART = (Inches(4.71), Inches(1.52), Inches(4.34), Inches(2.44))
_L4_LEFT_LABEL  = (Inches(0.48), Inches(1.32), Inches(3.29), Inches(0.27))
_L4_RIGHT_LABEL = (Inches(5.08), Inches(1.32), Inches(3.29), Inches(0.27))
_L4_LEFT_CAP    = (Inches(0.40), Inches(3.86), Inches(4.05), Inches(0.22))
_L4_RIGHT_CAP   = (Inches(5.57), Inches(3.82), Inches(4.05), Inches(0.22))


def _add_textbox(slide, ltwh, text, size, bold=False, color=DARK_GRAY):
    l, t, w, h = ltwh
    txb = slide.shapes.add_textbox(l, t, w, h)
    add_text(txb.text_frame, text, first=True, cn_font=CN_FONT, size=size, bold=bold, color=color)


def add_layout4_slide(prs, title, left_label='', right_label='',
                      left_caption='', right_caption=''):
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_slide_title(slide, title)
    if left_label:    _add_textbox(slide, _L4_LEFT_LABEL,  left_label,  10, color=DEEP_BLUE)
    if right_label:   _add_textbox(slide, _L4_RIGHT_LABEL, right_label, 10, color=DEEP_BLUE)
    if left_caption:  _add_textbox(slide, _L4_LEFT_CAP,    left_caption,  7, color=MID_GRAY)
    if right_caption: _add_textbox(slide, _L4_RIGHT_CAP,   right_caption, 7, color=MID_GRAY)
    return slide, _L4_LEFT_CHART, _L4_RIGHT_CHART

# ── 封底 ─────────────────────────────────────────────────────

def add_cn_closing_slide(prs, thanks='谢谢', contacts=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        tf  = ph.text_frame; tf.word_wrap = True
        if idx == 10:
            para = add_text(tf, thanks, first=True, cn_font=CN_FONT, en_font=None,
                            size=45, bold=True, color=DEEP_BLUE)
            _set_para_spacing(para, line_spc_pct=100)
        elif idx == 11 and contacts:
            for i, line in enumerate(contacts):
                para = add_text(tf, line, first=(i == 0), cn_font=CN_FONT, en_font=None,
                                size=12, bold=False, color=DEEP_BLUE)
                _set_para_spacing(para, line_spc_pct=170)
    return slide

# ═════════════════════════════════════════════════════════════
# 生成演示文稿
# ═════════════════════════════════════════════════════════════

prs = Presentation(TEMPLATE)

# 1. 提取目录表格 XML（必须在删除 slides 之前）
toc_table_xml = None
for slide in prs.slides:
    if slide.slide_layout.name == '中文目录页，仅供目录页使用':
        for shape in slide.shapes:
            if shape.shape_type == 19:
                toc_table_xml = copy.deepcopy(shape._element)
                break
        break

# 2. 删除所有原有幻灯片
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)

# ── 章节列表 ─────────────────────────────────────────────────
CHAPTERS = [
    "公司介绍",
    "投资策略与理念",
    "主要产品线",
    "业绩回顾",
    "风险管理",
    "展望与规划",
]

# ── Slide 1：中文封面 ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[0])
for shape in slide.shapes:
    ph = getattr(shape, 'placeholder_format', None)
    if ph is None: continue
    tf = shape.text_frame; tf.word_wrap = True
    if ph.idx == 1:
        add_text(tf, "易方达基金管理有限公司", first=True, size=28, bold=False)
        add_text(tf, "2024年度投资策略报告",   size=22, bold=True)
    elif ph.idx == 10:
        add_text(tf, "汇报人：张伟    投资研究部", first=True, size=14, bold=False)
        add_text(tf, "2024年3月·广州",           size=14, bold=False)

# ── Slide 2：目录（总览，无高亮） ────────────────────────────
add_toc_slide(prs, toc_table_xml, CHAPTERS, active_idx=-1)

# ── Slide 3：目录（当前章节：公司介绍） ──────────────────────
add_toc_slide(prs, toc_table_xml, CHAPTERS, active_idx=0)

# ── Slide 4：公司介绍（Layout 2 纯文字） ─────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_slide_title(slide, '公司介绍')
for shape in slide.shapes:
    try:
        if shape.placeholder_format.idx == 10:
            add_body_content(shape.text_frame, [
                ('公司概况', 0),
                ('成立于2001年，总部位于广州，注册资本15亿元', 1),
                ('截至2024年底，管理公募基金规模突破2万亿元，行业前三', 1),
                ('员工超过3000人，投研团队300余人', 1),
                ('主营业务', 0),
                ('公募基金：权益类、固收类、货币类、QDII', 1),
                ('专户理财：机构定制化资产管理', 1),
                ('基金投顾：为个人投资者提供组合配置服务', 1),
            ])
    except: pass

# ── Slide 5：目录（当前章节：投资策略） ──────────────────────
add_toc_slide(prs, toc_table_xml, CHAPTERS, active_idx=1)

# ── Slide 6：投资策略（Layout 3 左文右图） ───────────────────
slide = prs.slides.add_slide(prs.slide_layouts[3])
set_slide_title(slide, '投资策略与理念')
for shape in slide.shapes:
    try:
        if shape.placeholder_format.idx == 10:
            add_body_content(shape.text_frame, [
                ('核心理念', 0),
                ('客户至上：以客户长期利益为首要目标', 1),
                ('价值导向：基本面研究驱动投资决策', 1),
                ('长期视角：持有优质资产，穿越市场周期', 1),
                ('选股框架', 0),
                ('行业空间大、竞争格局好', 1),
                ('管理层优秀、公司治理规范', 1),
                ('估值合理或低估', 1),
            ])
    except: pass
# 右侧添加简单文本框代替图表（端到端示意）
txb = slide.shapes.add_textbox(Inches(5.4), Inches(1.42), Inches(4.34), Inches(2.71))
add_text(txb.text_frame, '[右侧图表区域]', first=True,
         cn_font=CN_FONT, size=10, bold=False, color=MID_GRAY)

# ── Slide 7：目录（当前章节：主要产品线） ────────────────────
add_toc_slide(prs, toc_table_xml, CHAPTERS, active_idx=2)

# ── Slide 8：主要产品线（Layout 4 双图） ─────────────────────
slide, l_area, r_area = add_layout4_slide(
    prs, '主要产品线规模分布',
    left_label='权益类基金规模（亿元）',
    right_label='固收类基金规模（亿元）',
    left_caption='数据来源：公司定期报告',
    right_caption='数据来源：公司定期报告',
)
# 左图
cd_left = CategoryChartData()
cd_left.categories = ['2021', '2022', '2023', '2024']
cd_left.add_series('权益类', (6800, 7200, 7900, 8600))
chart_l = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, *l_area, cd_left).chart
chart_l.has_title = False
chart_l.series[0].format.fill.solid()
chart_l.series[0].format.fill.fore_color.rgb = DEEP_BLUE

# 右图
cd_right = CategoryChartData()
cd_right.categories = ['2021', '2022', '2023', '2024']
cd_right.add_series('固收类', (5200, 6100, 7400, 8900))
chart_r = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, *r_area, cd_right).chart
chart_r.has_title = False
chart_r.series[0].format.fill.solid()
chart_r.series[0].format.fill.fore_color.rgb = BRIGHT_BLUE

# ── Slide 9：目录（当前章节：业绩回顾） ──────────────────────
add_toc_slide(prs, toc_table_xml, CHAPTERS, active_idx=3)

# ── Slide 10：业绩回顾（Layout 2） ───────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_slide_title(slide, '业绩回顾')
for shape in slide.shapes:
    try:
        if shape.placeholder_format.idx == 10:
            add_body_content(shape.text_frame, [
                ('权益类基金', 0),
                ('2024年旗舰主动权益基金平均超额收益+8.2%，跑赢基准', 1),
                ('易方达蓝筹精选混合：+22.3%，同类排名前5%', 1),
                ('易方达科技创新：+18.7%，重仓半导体与新能源', 1),
                ('固收类基金', 0),
                ('纯债基金平均收益4.1%，最大回撤控制在0.5%以内', 1),
                ('"固收+"产品平均收益6.8%，夏普比率优于同类均值', 1),
            ])
    except: pass

# ── Slide 11：中文封底 ────────────────────────────────────────
add_cn_closing_slide(prs,
    thanks='谢谢',
    contacts=[
        '联系我们：',
        '张伟    Tel：+86(20)8510-xxxx',
        'Email：investor@efunds.com.cn',
    ]
)

# ── 保存 ─────────────────────────────────────────────────────
out = os.path.join(SKILL_DIR, "test_samples", "e2e_test.pptx")
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
