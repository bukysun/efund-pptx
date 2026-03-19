"""Cover page verification test - generates 3 cover slide variants."""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree

TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "template.pptx")
CN_FONT  = "华文黑体_易方达"
EN_FONT  = "Arial"
DEEP_BLUE = RGBColor(0, 80, 150)

def qn(tag):
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    prefix, local = tag.split(':')
    return f'{{{nsmap[prefix]}}}{local}'

def apply_font(run, cn_font=CN_FONT, en_font=None, size=None, bold=None, color=None):
    rPr = run._r.get_or_add_rPr()
    if cn_font and not en_font:
        latin = rPr.find(qn('a:latin'))
        if latin is None: latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', cn_font); latin.set('charset', '-122')
        cs = rPr.find(qn('a:cs'))
        if cs is None: cs = etree.SubElement(rPr, qn('a:cs'))
        cs.set('typeface', cn_font); cs.set('charset', '-122')
    elif en_font and not cn_font:
        latin = rPr.find(qn('a:latin'))
        if latin is None: latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', en_font); latin.attrib.pop('charset', None)
    if size  is not None: run.font.size      = Pt(size)
    if bold  is not None: run.font.bold      = bold
    if color is not None: run.font.color.rgb = color

def add_text(tf, text, first=False, cn_font=CN_FONT, en_font=None,
             size=None, bold=None, color=DEEP_BLUE):
    if first:
        para = tf.paragraphs[0]; para.clear()
    else:
        para = tf.add_paragraph()
    run = para.add_run(); run.text = text
    apply_font(run, cn_font=cn_font, en_font=en_font, size=size, bold=bold, color=color)
    return para

prs = Presentation(TEMPLATE)

# 删除模板中已有的所有幻灯片，只保留 slide layouts
from pptx.oxml.ns import qn as _qn
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)

# ── 中文封面（Layout 0）─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[0])
for shape in slide.shapes:
    ph = getattr(shape, 'placeholder_format', None)
    if ph is None: continue
    tf = shape.text_frame; tf.word_wrap = True
    if ph.idx == 1:
        add_text(tf, "易方达的历史和文化", first=True, size=28, bold=False)
        add_text(tf, "2024年新员工培训",   size=22, bold=True)
    elif ph.idx == 10:
        add_text(tf, "汇报人：张三", first=True, size=14, bold=False)
        add_text(tf, "2024年3月",    size=14, bold=False)

# ── 英文封面（Layout 7，纯英文）────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[7])
for shape in slide.shapes:
    ph = getattr(shape, 'placeholder_format', None)
    if ph is None: continue
    tf = shape.text_frame; tf.word_wrap = True
    if ph.idx == 1:
        add_text(tf, "E Fund Annual Report 2024",
                 first=True, cn_font=None, en_font=EN_FONT, size=28, bold=False)
        add_text(tf, "Strategic Overview",
                 cn_font=None, en_font=EN_FONT, size=21, bold=True)
    elif ph.idx == 10:
        add_text(tf, "Presenter: John Smith",
                 first=True, cn_font=None, en_font=EN_FONT, size=14, bold=False)
        add_text(tf, "March 2024, Shanghai",
                 cn_font=None, en_font=EN_FONT, size=14, bold=False)

# ── 中英文封面（Layout 7，首行中文）────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[7])
for shape in slide.shapes:
    ph = getattr(shape, 'placeholder_format', None)
    if ph is None: continue
    tf = shape.text_frame; tf.word_wrap = True
    if ph.idx == 1:
        add_text(tf, "易方达年度报告 2024", first=True, size=28, bold=False)
        add_text(tf, "Annual Report",
                 cn_font=None, en_font=EN_FONT, size=21, bold=True)
    elif ph.idx == 10:
        add_text(tf, "汇报人：李四", first=True, size=14, bold=False)
        add_text(tf, "Presenter: Li Si",
                 cn_font=None, en_font=EN_FONT, size=12, bold=False)
        add_text(tf, "March 2024",
                 cn_font=None, en_font=EN_FONT, size=12, bold=False)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "test_samples", "cover_test.pptx")
prs.save(out)
print(f"Saved → {out}")
